"""Dashboard export builders (Block E.3): management-ready PDF and PPTX."""

import io
import logging
from datetime import datetime, timezone
from typing import Any, Dict, List

logger = logging.getLogger(__name__)


def _period_label(data: Dict[str, Any]) -> str:
    return f"последние {data.get('days', 30)} дней"


def build_dashboard_pdf(data: Dict[str, Any]) -> bytes:
    """Multi-section dashboard report (reuses the Block D PDF stack)."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    styles = getSampleStyleSheet()
    h2 = ParagraphStyle("h2", parent=styles["Heading2"], fontSize=13, spaceBefore=10)
    body = ParagraphStyle("body", parent=styles["BodyText"], fontSize=9)
    note = ParagraphStyle("note", parent=body, fontSize=7.5,
                          textColor=colors.HexColor("#555555"))

    def table(rows: List[List], widths) -> Table:
        t = Table(rows, colWidths=widths)
        t.setStyle(TableStyle([
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
            ("GRID", (0, 0), (-1, -1), 0.3, colors.grey),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ]))
        return t

    story = [
        Paragraph("FakeDetect — Отчёт по защите бренда", styles["Title"]),
        Paragraph(
            f"Период: {_period_label(data)} · "
            f"Сгенерировано (UTC): {datetime.now(timezone.utc).isoformat()}",
            body,
        ),
        Spacer(1, 5 * mm),
    ]

    s = data.get("summary", {})
    story.append(Paragraph("1. Ключевые показатели", h2))
    story.append(table(
        [["Всего проверок", "Подделок", "Подозрительных", "Оригиналов"],
         [str(s.get("total", 0)), str(s.get("fakes", 0)),
          str(s.get("suspicious", 0)), str(s.get("originals", 0))]],
        [40 * mm] * 4,
    ))

    ts = data.get("timeseries", [])
    if ts:
        story.append(Paragraph(f"2. Динамика вердиктов ({data.get('granularity')})", h2))
        rows = [["Период", "Проверок", "Подделок", "Подозрительных", "Оригиналов"]]
        rows += [[p["period"], p["total"], p["fakes"], p["suspicious"],
                  p["originals"]] for p in ts[-20:]]
        story.append(table(rows, [30 * mm, 25 * mm, 25 * mm, 32 * mm, 28 * mm]))

    sellers = data.get("top_sellers", [])
    if sellers:
        story.append(Paragraph("3. Топ продавцов-нарушителей", h2))
        rows = [["Продавец", "Проверок", "Нарушений", "Подделок", "Ср. уверенность"]]
        rows += [[sv["seller"], sv["total_checks"], sv["violations"],
                  sv["fakes"], f"{sv['avg_confidence']}%"] for sv in sellers[:10]]
        story.append(table(rows, [50 * mm, 22 * mm, 25 * mm, 22 * mm, 28 * mm]))

    rev = data.get("protected_revenue", {})
    story.append(Paragraph("4. Оценка защищённой выручки", h2))
    est = rev.get("protected_revenue_estimate")
    story.append(Paragraph(
        f"Подтверждённых подделок за период: {rev.get('confirmed_fakes', 0)}. "
        f"Средняя цена оригинала: {rev.get('avg_original_price') or '—'} ₽. "
        f"<b>Оценка защищённой выручки: "
        f"{f'{est:,} ₽'.replace(',', ' ') if est is not None else '—'}</b>",
        body,
    ))
    story.append(Paragraph(rev.get("disclaimer", ""), note))

    timing = data.get("timing", {})
    story.append(Paragraph("5. Операционные метрики", h2))
    story.append(Paragraph(
        f"Time-to-detection: {timing.get('time_to_detection_days') or '—'} дн. · "
        f"Time-to-resolution: {timing.get('time_to_resolution_days') or '—'} дн.",
        body,
    ))
    story.append(Paragraph(timing.get("note", ""), note))

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4,
                            leftMargin=15 * mm, rightMargin=15 * mm,
                            topMargin=12 * mm, bottomMargin=12 * mm)
    doc.build(story)
    return buffer.getvalue()


def build_dashboard_pptx(data: Dict[str, Any]) -> bytes:
    """Compact management deck: title, key numbers, top violators."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "FakeDetect — защита бренда"
    slide.placeholders[1].text = (
        f"Отчёт за {_period_label(data)} · "
        f"{datetime.now(timezone.utc).strftime('%Y-%m-%d')}"
    )

    s = data.get("summary", {})
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Ключевые показатели"
    tf = slide.placeholders[1].text_frame
    bullets = [
        f"Проверок карточек: {s.get('total', 0)}",
        f"Подтверждённых подделок: {s.get('fakes', 0)}",
        f"Подозрительных: {s.get('suspicious', 0)}",
        f"Оригиналов: {s.get('originals', 0)}",
    ]
    rev = data.get("protected_revenue", {})
    revenue_disclaimer = None
    if rev.get("protected_revenue_estimate") is not None:
        est = rev["protected_revenue_estimate"]
        bullets.append(f"Оценка защищённой выручки: {est:,} ₽".replace(",", " "))
        # E-C4: the estimate must never appear on a slide without its caveat —
        # a CFO reading only the bullet would otherwise take it as an exact figure.
        revenue_disclaimer = rev.get("disclaimer")
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(20)
    if revenue_disclaimer:
        p = tf.add_paragraph()
        p.text = f"* {revenue_disclaimer}"
        p.font.size = Pt(11)
        p.font.italic = True

    sellers = data.get("top_sellers", [])
    if sellers:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Топ продавцов-нарушителей"
        tf = slide.placeholders[1].text_frame
        for i, sv in enumerate(sellers[:8]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = (f"{sv['seller']}: подделок {sv['fakes']}, "
                      f"нарушений {sv['violations']}")
            p.font.size = Pt(18)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
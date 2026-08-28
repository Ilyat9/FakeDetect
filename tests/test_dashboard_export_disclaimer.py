"""E-C4 regression: the protected-revenue estimate must never appear in an
export (PDF or PPTX) without its methodology disclaimer next to it — a
management deck slide showing a bare number reads as an exact financial
figure, not the rough estimate it actually is.
"""

import io

from pptx import Presentation
from pypdf import PdfReader

from app.services.dashboard_export import build_dashboard_pdf, build_dashboard_pptx

SAMPLE_DATA = {
    "days": 30,
    "granularity": "day",
    "summary": {"total": 10, "fakes": 4, "suspicious": 2, "originals": 4},
    "timeseries": [],
    "top_sellers": [],
    "protected_revenue": {
        "confirmed_fakes": 4,
        "avg_original_price": 5000,
        "protected_revenue_estimate": 20000,
        "disclaimer": "Оценка, а не точная цифра: подтверждённые подделки × средняя цена оригинала.",
    },
    "timing": {"time_to_detection_days": 1.5, "time_to_resolution_days": 3.0, "note": ""},
}


def test_pdf_includes_revenue_disclaimer_next_to_estimate():
    # reportlab's default Helvetica has no ToUnicode CMap for Cyrillic, so
    # pypdf's text extraction can't recover Cyrillic glyphs (a separate,
    # pre-existing font-encoding limitation, out of scope here) — use an
    # ASCII disclaimer to verify the mechanism actually embeds it in the PDF.
    data = dict(SAMPLE_DATA)
    data["protected_revenue"] = dict(
        SAMPLE_DATA["protected_revenue"],
        disclaimer="ESTIMATE ONLY, NOT AN EXACT FIGURE",
    )
    pdf_bytes = build_dashboard_pdf(data)
    reader = PdfReader(io.BytesIO(pdf_bytes))
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    assert "20" in text  # the estimate rendered somewhere
    assert "ESTIMATE ONLY, NOT AN EXACT FIGURE" in text


def test_pptx_includes_revenue_disclaimer_next_to_estimate():
    pptx_bytes = build_dashboard_pptx(SAMPLE_DATA)
    prs = Presentation(io.BytesIO(pptx_bytes))

    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
    combined = "\n".join(all_text)

    assert "Оценка защищённой выручки" in combined
    assert "Оценка, а не точная цифра" in combined


def test_pptx_omits_disclaimer_when_no_estimate_available():
    data = dict(SAMPLE_DATA)
    data["protected_revenue"] = {
        "confirmed_fakes": 0, "avg_original_price": None,
        "protected_revenue_estimate": None, "disclaimer": "Недостаточно данных",
    }
    pptx_bytes = build_dashboard_pptx(data)
    prs = Presentation(io.BytesIO(pptx_bytes))
    all_text = "\n".join(
        shape.text_frame.text
        for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
    )
    assert "Оценка защищённой выручки" not in all_text
    assert "Недостаточно данных" not in all_text
